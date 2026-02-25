from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Colors
    RED = RGBColor(211, 47, 47)     # #D32F2F
    BLACK = RGBColor(33, 33, 33)    # #212121
    WHITE = RGBColor(255, 255, 255) # #FFFFFF
    GRAY = RGBColor(117, 117, 117)  # #757575

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        title.text_frame.radians = 0
        
        # Style Title
        p = title.text_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER

        subtitle.text = subtitle_text
        p = subtitle.text_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER

    def add_content_slide(title_text, content_items):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        p = title.text_frame.paragraphs[0]
        p.font.color.rgb = RED
        p.font.bold = True
        
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear() # Clear existing dummy text

        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_after = Pt(10)
            if item.startswith("  -"):
                p.level = 1
                p.text = item.replace("  -", "").strip()
    
    def add_section_header(title_text):
        slide_layout = prs.slide_layouts[2] # Section Header
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        p = title.text_frame.paragraphs[0]
        p.font.color.rgb = WHITE
        
        # Background color hack (simple version)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RED

    # --- 1. Portada ---
    add_title_slide("LEVEL UP", "ORDEN. CONTINUIDAD. CONTROL. RESULTADOS.")

    # --- 2. Introducción ---
    add_content_slide("Introducción: Quién es Level Up", [
        "Level Up innova y ayuda a los negocios a pasar al siguiente nivel.",
        "Implementamos estrategias integrales en servicio al cliente telefónico.",
        "Optimizamos la atención de llamadas y mejoramos la experiencia del cliente.",
        "Resultados: Mayor control, menos estrés operativo, más ventas.",
        "Enfoque en crecer y destacar frente a la competencia."
    ])

    # --- 3. Problemática Actual ---
    add_content_slide("Problemática Actual", [
        "Alta demanda de atención telefónica:",
        "  - Largas esperas y pérdida de ventas.",
        "  - Atención inconsistente que daña la reputación.",
        "Logística de entrega deficiente:",
        "  - Reparto lento y riesgos físicos (accidentes, pérdidas).",
        "  - Costos por devoluciones.",
        "Dependencia de Apps externas:",
        "  - Altas comisiones y pérdida de control del cliente."
    ])

    # --- 4. Impacto ---
    add_content_slide("Impacto en el Negocio", [
        "Incapacidad de atender a todos los clientes eficazmente.",
        "Pérdida de clientes habituales y nuevos.",
        "Aumento de costos operativos y reducción de rentabilidad.",
        "La atención telefónica es crítica para la fidelización."
    ])

    # --- 5. La Solución Level Up ---
    add_content_slide("La Solución Level Up", [
        "Centralizamos y organizamos la atención telefónica.",
        "Garantizamos respuesta a todas las llamadas.",
        "Captura completa de pedidos y datos.",
        "Envío directo de información al negocio.",
        "“Le quitamos el celular de las manos al negocio para que pueda enfocarse en lo importante.”"
    ])

    # --- 6. Modelo de Trabajo ---
    add_content_slide("Modelo de Trabajo", [
        "Implementación: Configuración de desvío o línea dedicada.",
        "Atención: Equipo profesional con protocolos.",
        "Gestión: Captura de orden y detalles.",
        "Transmisión: Envío procesado listo para preparar.",
        "Beneficio: Libera tiempo y agiliza entregas."
    ])

    # --- 7. Paquetes: Level 1 ---
    add_content_slide("Paquete Level 1", [
        "Precio: $1,000 MXN / Mensual",
        "Promoción: 1 semana gratis",
        "Incluye:",
        "  - Atención total de llamadas",
        "  - Captura de pedidos",
        "  - Envío automático de comandas",
        "  - Eliminación de la carga del celular"
    ])

    # --- 8. Paquetes: Level 2 ---
    add_content_slide("Paquete Level 2", [
        "Precio: $2,000 MXN / Mensual",
        "Incluye todo lo de Level 1, más:",
        "  - Incremento del tráfico de llamadas",
        "  - Publicidad y posicionamiento",
        "  - Fotos, videos y publicaciones en redes sociales",
        "Ideal para negocios que buscan más llamadas."
    ])

    # --- 9. Paquetes: Level 3 ---
    add_content_slide("Paquete Level 3", [
        "Opción A: $20,000 MXN / Mensual",
        "Opción B: $70 MXN / Envío",
        "Incluye todo lo de Level 1 y 2, más:",
        "  - Promociones exclusivas por llamada",
        "  - Servicio de entrega con DRON exclusivo Level Up",
        "  - Dron proporcionado al firmar contrato",
        "  - Instalación de área de aterrizaje"
    ])

    # --- 10. Misión y Visión ---
    add_content_slide("Misión y Visión", [
        "Misión:",
        "  - Ayudar a negocios a mantener control total.",
        "  - Atención ordenada, continua y profesional.",
        "  - Crecimiento estable y sostenible.",
        "Visión:",
        "  - Nuevo estándar en gestión de atención al cliente.",
        "  - Orden, continuidad y control como base del crecimiento."
    ])

    # --- 11. Valores ---
    add_content_slide("Valores Level Up", [
        "El cliente siempre es atendido.",
        "Simplicidad operativa.",
        "El tiempo vale más que el dinero.",
        "Resultados medibles.",
        "Innovación con propósito.",
        "Crecimiento compartido."
    ])

    # --- 12. Manifiesto ---
    add_content_slide("Manifiesto de Marca", [
        "ORDEN sobre el caos.",
        "CONTINUIDAD de la excelencia.",
        "CONTROL que libera.",
        "ATENCIÓN CONSTANTE.",
        "CRECIMIENTO sin caos."
    ])

    # --- 13. Cierre ---
    add_title_slide("Gracias", "¿Listo para el siguiente nivel?\nContáctanos hoy.")

    prs.save("LEVEL_UP_Presentacion.pptx")
    print("Presentación guardada como LEVEL_UP_Presentacion.pptx")

if __name__ == "__main__":
    create_presentation()
