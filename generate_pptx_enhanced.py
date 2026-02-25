from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_enhanced_presentation():
    pres = Presentation()

    # Define Colors
    GUINDA = RGBColor(157, 36, 73)
    VERDE = RGBColor(40, 92, 77)
    NARANJA = RGBColor(179, 93, 0)
    GRIS_OSCURO = RGBColor(51, 51, 51)

    def add_title_slide(pres, title_text, subtitle_text):
        slide_layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text
        for par in title.text_frame.paragraphs:
            par.font.color.rgb = GUINDA
            par.font.bold = True
            par.font.size = Pt(40)
        for par in subtitle.text_frame.paragraphs:
            par.font.size = Pt(22)

    def add_bullet_slide(pres, title_text, bullets, title_color=GUINDA):
        slide_layout = pres.slide_layouts[1]
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        for par in title.text_frame.paragraphs:
            par.font.color.rgb = title_color
            par.font.bold = True
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.word_wrap = True
        for i, point in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(20)
            p.space_after = Pt(10)

    # 1. Cover
    add_title_slide(pres, "Control del Sarampión: Tres Indicadores Clave", 
                    "Estrategia Nacional de Vacunación\nSeguimiento al 13 de febrero de 2026\nSecretaría de Salud | México")

    # 2. Panorama
    add_bullet_slide(pres, "Panorama Epidemiológico Actual", [
        "Casos Totales Confirmados: 9,712",
        "Casos Activos: 1,185",
        "Defunciones: 29",
        "Tasa de Letalidad Nacional: 0.30%",
        "24 de 32 entidades con baja incidencia acumulada (<100 casos)."
    ])

    # 3. Los 3 Indicadores
    add_bullet_slide(pres, "Los Tres Pilares del Monitoreo", [
        "Indicador 1: Inmunización Temprana (6 a 11 meses).",
        "Indicador 2: Consolidación Infantil (19 meses a 12 años).",
        "Indicador 3: Recuperación Adulta (13 a 49 años).",
        "Meta Homogénea Nacional: Cobertura ≥ 95%."
    ], title_color=VERDE)

    # 4. Indicador 1
    add_bullet_slide(pres, "Indicador 1: Inmunización Temprana", [
        "Población: Lactantes de 6 a 11 meses (Dosis 0).",
        "Cobertura Actual: 67%",
        "Meta: 95%",
        "Brecha a cerrar: 28% puntos porcentuales.",
        "Prioridad en entidades con casos activos."
    ])

    # 5. Indicador 2
    add_bullet_slide(pres, "Indicador 2: Consolidación Infantil", [
        "Población: Niños y niñas de 19 meses a 12 años.",
        "Cobertura Actual: 81%",
        "Identificado como el grupo más vulnerable.",
        "Meta: Alcanzar el 95% para inmunidad de rebaño.",
        "Brecha a cerrar: 14% puntos porcentuales."
    ], title_color=VERDE)

    # 6. Indicador 3
    add_bullet_slide(pres, "Indicador 3: Recuperación Adulta", [
        "Población: Adultos de 13 a 49 años.",
        "Cobertura Actual: 42%",
        "Foco: 7 Entidades Federativas con alta incidencia.",
        "Meta de recuperación de coberturas históricas.",
        "Uso de vacuna Doble Viral (SR)."
    ], title_color=NARANJA)

    # 7. Metas Institucionales
    add_bullet_slide(pres, "Metas de 1ra Dosis por Institución", [
        "Meta Nacional Total: 638,647 dosis.",
        "Secretaría de Salud (SSA): 198,990 dosis.",
        "IMSS-BIENESTAR: 245,134 dosis.",
        "IMSS-ORDINARIO: 150,695 dosis.",
        "Esfuerzo coordinado para lograr el 95% antes de noviembre."
    ])

    # 8. Directrices
    add_bullet_slide(pres, "Directrices de Acción", [
        "Menores de 12 años: Acudir de inmediato si no tienen esquema completo.",
        "Aplicar refuerzo a los 6 meses de la primera dosis.",
        "13 a 49 años: Vacunación prioritaria en centros de salud y brigadas.",
        "Uso de plataforma nacional para validación de disponibilidad."
    ])

    # 9. Logística
    add_bullet_slide(pres, "Logística y Garantía de Abasto", [
        "Dosis Distribuidas en Estados: 16.7 Millones.",
        "Dosis disponibles para envío inmediato: 11.0 Millones.",
        "Suficiencia total para cubrir las metas nacionales.",
        "Distribución diaria basada en reporte de inventarios almacenes."
    ], title_color=VERDE)

    # 10. Ruta de Éxito
    add_bullet_slide(pres, "Ruta de Éxito en Territorio", [
        "Seguimiento diario de la evolución de casos por municipio.",
        "Garantizar abasto en los 21,000 puntos de vacunación.",
        "Actualización diaria de plataforma sectorial.",
        "Nombramiento de enlaces exclusivos para monitoreo."
    ])

    # 11. Contacto
    add_bullet_slide(pres, "Consulta de Puntos de Vacunación", [
        "Más de 21,000 puntos disponibles a nivel nacional.",
        "Portal Web: dondemevacuno.salud.gob.mx",
        "¡Tu protección es nuestra prioridad!",
        "Acude hoy mismo."
    ])

    # Save
    path = "C:\\Users\\aicil\\OneDrive\\Escritorio\\PVU\\SARAMPIÓN\\BOLETINES\\Estrategia_Indicadores_Sarampion_2026_Completa.pptx"
    pres.save(path)
    print(f"Enhanced PPTX saved to: {path}")

if __name__ == "__main__":
    create_enhanced_presentation()
