from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def update_presentation():
    pptx_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\Estrategia_Indicadores_Sarampion_2026_Completa16FEBRERO2026.pptx"
    output_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\Estrategia_Indicadores_Sarampion_2026_Actualizada_17022026.pptx"
    
    prs = Presentation(pptx_path)

    # 1. Update Title Slide
    title_slide = prs.slides[0]
    for shape in title_slide.shapes:
        if shape.has_text_frame:
            text = shape.text.upper()
            if "CONTROL DEL SARAMPIÓN" in text or "ESTRATEGIA NACIONAL" in text or "SEGUIMIENTO" in text:
                # Assuming top shape is title, middle is subtitle/date
                if shape.top < Inches(2): # Likely Title
                    shape.text_frame.text = "4ª REUNIÓN EXTRAORDINARIA DEL COEVA\nSITUACIÓN EPIDEMIOLÓGICA DEL SARAMPIÓN"
                else:
                    shape.text_frame.text = "Martes 17 de febrero de 2026 | 08:15 Horas\nModalidad Virtual (Zoom)"

    # 2. Insert Agenda Slide (Slide 2)
    # Using layout 1 (Title and Content)
    slide_layout = prs.slide_layouts[1]
    agenda_slide = prs.slides.add_slide(slide_layout)
    # Move agenda slide to position 2
    prs.slides._sldIdLst.insert(1, prs.slides._sldIdLst[-1])
    
    title_shape = agenda_slide.shapes.title
    title_shape.text = "Orden del Día"
    
    body_shape = agenda_slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "1. Lista de asistencia y declaratoria de quórum."
    p = tf.add_paragraph(); p.text = "2. Lectura y aprobación del orden del día."
    p = tf.add_paragraph(); p.text = "3. Panorama epidemiológico estatal (Casos confirmados: 333)."
    p = tf.add_paragraph(); p.text = "4. Nuevas directrices nacionales (Priorización 6m-12a)."
    p = tf.add_paragraph(); p.text = "5. Avance operativo de vacunación por jurisdicción."
    p = tf.add_paragraph(); p.text = "6. Monitoreo de existencias SRP/SR e insumos."
    p = tf.add_paragraph(); p.text = "7. Seguimiento de acuerdos y acciones prioritarias."
    p = tf.add_paragraph(); p.text = "8. Asuntos generales."
    p = tf.add_paragraph(); p.text = "9. Cierre de sesión."

    # 3. Update Panorama Epidemiológico (formerly Slide 2, now Slide 3)
    panorama_slide = prs.slides[2]
    # Update title and subtitle for Durango data
    for shape in panorama_slide.shapes:
        if shape.has_text_frame:
            if "Panorama Epidemiológico Actual" in shape.text:
                shape.text_frame.text = "Panorama Epidemiológico Estatal"
            elif "Casos Totales Confirmados" in shape.text or "Casos Activos" in shape.text:
                shape.text_frame.text = "Casos Confirmados en Durango: 333\nFocos Activos: Victoria de Durango, Mezquital\n(Corte 16 Feb)"

    # 4. New Directrices Slide (at position 4)
    directrices_slide = prs.slides.add_slide(prs.slide_layouts[1])
    prs.slides._sldIdLst.insert(3, prs.slides._sldIdLst[-1]) 
    directrices_slide.shapes.title.text = "Nuevas Directrices Nacionales"
    tf = directrices_slide.shapes.placeholders[1].text_frame
    tf.text = "• Priorización: Niñas y niños de 6 meses a 12 años."
    p = tf.add_paragraph(); p.text = "• Esquema Completo: No refuerzos innecesarios si ya tiene 2 dosis."
    p = tf.add_paragraph(); p.text = "• Logística: Distribución inmediata a puntos de vacunación."
    p = tf.add_paragraph(); p.text = "• Reporte: Actualización diaria de plataforma y existencias."
    p = tf.add_paragraph(); p.text = "• Enlace: Ratificación de enlace estatal de monitoreo."
    p = tf.add_paragraph(); p.text = "• Gobernanza: Reuniones semanales de alto nivel."

    # 5. Update Vaccination Progress (formerly Slide 7, now likely Slide 9)
    # Finding slide by content
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "Metas de 1ra Dosis" in shape.text:
                shape.text_frame.text = "Avance Operativo de Vacunación"
                # Add bullet points for strategies
                content_shape = slide.shapes.placeholders[1]
                ctf = content_shape.text_frame
                ctf.text = "• Dosis aplicadas según Jurisdicción (Corte 13 Feb)."
                p = ctf.add_paragraph(); p.text = "• Estrategias Activas: Barrido, Brigadas Móviles, Escolar."
                p = ctf.add_paragraph(); p.text = "• Meta: Cobertura ≥ 95% en grupo prioritario."

    # 6. Update Inventory & Logistics (formerly Slide 9, now Slide 11)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "Logística y Garantía de Abasto" in shape.text:
                content_shape = slide.shapes.placeholders[1]
                ctf = content_shape.text_frame
                ctf.text = "• Monitoreo de existencias SRP/SR e insumos críticos."
                p = ctf.add_paragraph(); p.text = "• Dosis distribuidas en Estados: 16.7 Millones."
                p = ctf.add_paragraph(); p.text = "• Acciones inmediatas de redistribución de biológico."
                p = ctf.add_paragraph(); p.text = "• Registro diario obligatorio en plataforma nacional."

    # 7. Add "Acuerdos y Seguimiento" Slide at the end
    acuerdos_slide = prs.slides.add_slide(prs.slide_layouts[1])
    acuerdos_slide.shapes.title.text = "Acuerdos y Acciones Prioritarias"
    tf = acuerdos_slide.shapes.placeholders[1].text_frame
    tf.text = "• Garantizar la distribución oportuna de biológico a centros de salud."
    p = tf.add_paragraph(); p.text = "• Reforzar el monitoreo diario de existencias y dosis aplicadas."
    p = tf.add_paragraph(); p.text = "• Intensificar la vacunación en población de 6 meses a 12 años."
    p = tf.add_paragraph(); p.text = "• Próxima reunión de seguimiento: Martes 24 de febrero."

    # Save the updated presentation
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    update_presentation()
