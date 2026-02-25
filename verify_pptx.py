from pptx import Presentation
import os

def verify_pptx():
    pptx_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\Estrategia_Indicadores_Sarampion_2026_Actualizada_17022026.pptx"
    
    if not os.path.exists(pptx_path):
        print(f"Error: {pptx_path} does not exist.")
        return

    prs = Presentation(pptx_path)
    print(f"Total Slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"  - {shape.text[:100]}...")
    
    print("\nVerification complete.")

if __name__ == "__main__":
    verify_pptx()
