from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    pres = Presentation()

    # Define Colors (Mexican Health Institutional Colors)
    GUINDA = RGBColor(157, 36, 73)
    VERDE = RGBColor(40, 92, 77)
    GRIS_OSCURO = RGBColor(51, 51, 51)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(pres, title_text, subtitle_text):
        slide_layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = subtitle_text
        
        # Style Title
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.color.rgb = GUINDA
            paragraph.font.bold = True
            paragraph.font.size = Pt(44)
            
        # Style Subtitle
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.color.rgb = GRIS_OSCURO
            paragraph.font.size = Pt(24)

    def add_content_slide(pres, title_text, bullet_points):
        slide_layout = pres.slide_layouts[1]
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = title_text
        # Style Title
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.color.rgb = GUINDA
            paragraph.font.bold = True
            
        tf = content.text_frame
        tf.text = bullet_points[0]
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = GRIS_OSCURO

    # --- SLIDES ---

    # 1. Cover
    add_title_slide(pres, "Estrategia Nacional de Atención al Sarampión", 
                    "Directrices de Vacunación y Seguimiento\nFebrero 2026\nGobierno de México | Salud")

    # 2. Panorama
    add_content_slide(pres, "Panorama Epidemiológico (Corte 13 Feb)", [
        "Casos Totales Confirmados: 9,712",
        "Casos Activos: 1,185",
        "Tasa de Letalidad: 0.30%",
        "Distribución: 24 de 32 entidades con <100 casos acumulados.",
        "Situación bajo control mediante vigilancia activa."
    ])

    # 3. El Pilar: Vacunación
    add_content_slide(pres, "La Clave: Vacunación Universal", [
        "Protección Directa: Evita formas graves y complicaciones.",
        "Control de Transmisión: Detiene la propagación comunitaria.",
        "Hoy contamos con disponibilidad nacional de biológicos.",
        "Garantizamos el acceso gratuito en todo el territorio."
    ])

    # 4. Grupo Prioritario: Niños
    add_content_slide(pres, "Prioridad: Menores de 12 Años", [
        "Población Objetivo: Niños de 6 meses a 12 años.",
        "Sin Vacuna: Acudir de inmediato a centro de salud.",
        "Esquema Incompleto: Aplicar refuerzo si pasaron 6 meses de la 1ª dosis.",
        "Esquema Completo (2 dosis): El menor ya está protegido."
    ])

    # 5. Población 13 a 49 años
    add_content_slide(pres, "Recuperación: Población 13 a 49 Años", [
        "Estados Prioritarios: Jalisco, Colima, Chiapas, Sinaloa, Nayarit, Tabasco y CDMX.",
        "Vacuna Doble Viral (SR) para quienes no tengan esquema completo.",
        "Resto del país: Inicio gradual a partir del mes de mayo.",
        "Objetivo: Cerrar brechas de inmunidad en adultos."
    ])

    # 6. Cobertura Actual
    add_content_slide(pres, "Estado Actual de Cobertura (Sectorial)", [
        "6 a 11 meses (Dosis 0): 67% de cobertura.",
        "19 meses a 12 años: 81% de cobertura.",
        "13 a 49 años: 42% de cobertura (en proceso de expansión).",
        "Meta: Superar el 95% en todos los grupos."
    ])

    # 7. Abasto y Distribución
    add_content_slide(pres, "Logística de Abasto", [
        "Dosis Distribuidas: 16.7 Millones ya en los estados.",
        "Dosis en Reserva/Proceso: 11 Millones adicionales listas.",
        "Total Disponibilidad: Suficiente para la Estrategia Nacional.",
        "Distribución diaria garantizada a centros de salud."
    ])

    # 8. Pasos Críticos para el Éxito
    add_content_slide(pres, "Pasos Críticos de la Estrategia", [
        "1. Cumplimiento de vacunación según grupos prioritarios.",
        "2. Distribución eficiente a nivel local (centros de salud).",
        "3. Monitoreo diario por cada entidad federativa.",
        "4. Reuniones semanales de alto nivel para toma de decisiones."
    ])

    # 9. Contacto
    add_content_slide(pres, "¿Dónde Vacunarse?", [
        "Consulta más de 21,000 puntos disponibles.",
        "Sitio Web Oficial: dondemevacuno.salud.gob.mx",
        "¡La prevención es nuestra mejor protección!",
        "Acude a tu unidad de salud más cercana."
    ])

    # Save
    file_path = "C:\\Users\\aicil\\OneDrive\\Escritorio\\PVU\\SARAMPIÓN\\BOLETINES\\Estrategia_Vacunacion_Sarampion_2026.pptx"
    pres.save(file_path)
    print(f"Presentation saved to: {file_path}")

if __name__ == "__main__":
    create_presentation()
