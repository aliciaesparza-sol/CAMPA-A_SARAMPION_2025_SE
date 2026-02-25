import os
from pypdf import PdfReader
from pptx import Presentation

def extract_pdf_text(filepath):
    try:
        reader = PdfReader(filepath)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error reading {filepath}: {e}"

def analyze_pptx(filepath):
    try:
        prs = Presentation(filepath)
        slides_info = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_info.append({"index": i + 1, "content": " | ".join(slide_text)})
        return slides_info
    except Exception as e:
        return f"Error reading {filepath}: {e}"

# Paths
pdf_gobernadores = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\BOLETINES\Sarampion con Gobernadores 2026 02 11 Estrategia menor 12 años.pdf"
pdf_seguimiento = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\BOLETINES\SEGUIMIENTO SARAMPION 13FEBRERO2026 ESTRATEGIA INDICADORES.pdf"
pdf_casos = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\CASOS NOTIFICADOS\CASOS CONFIRMADOS 2026\LISTADO CASOS CONFIRMADOS BROTE_SARAMPIÓN 16.02.2026.pdf"
pptx_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\Estrategia_Indicadores_Sarampion_2026_Completa16FEBRERO2026.pptx"

results = {}

# Extract and process
results["gobernadores"] = extract_pdf_text(pdf_gobernadores)
results["seguimiento"] = extract_pdf_text(pdf_seguimiento)
results["casos"] = extract_pdf_text(pdf_casos)
results["pptx"] = analyze_pptx(pptx_path)

# Output summary to a file with utf-8
with open("summary_data.txt", "w", encoding="utf-8") as f:
    f.write("=== GOBERNADORES / DIRECTRICES ===\n")
    f.write(results["gobernadores"][:2000] + "...\n\n")
    f.write("=== SEGUIMIENTO / INDICADORES ===\n")
    f.write(results["seguimiento"][:2000] + "...\n\n")
    f.write("=== CASOS CONFIRMADOS ===\n")
    # Count confirmed cases if possible
    lines = results["casos"].split("\n")
    f.write(f"Total lines in cases PDF: {len(lines)}\n")
    f.write(results["casos"][:1000] + "...\n\n")
    f.write("=== PPTX STRUCTURE ===\n")
    for s in results["pptx"]:
        f.write(f"Slide {s['index']}: {s['content'][:300]}\n")

print("Done. Check summary_data.txt")
